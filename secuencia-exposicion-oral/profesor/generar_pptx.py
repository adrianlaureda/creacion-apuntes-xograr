#!/usr/bin/env python3
"""
Genera PowerPoint de la presentación "Mi vida" con menú de navegación lateral.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Fuente
FONT_NAME = 'Helvetica Neue'

# Colores modo OSCURO
COLORS_DARK = {
    'bg_primary': RGBColor(0x05, 0x05, 0x05),
    'bg_card': RGBColor(0x11, 0x11, 0x11),
    'text_primary': RGBColor(0xFF, 0xFF, 0xFF),
    'text_secondary': RGBColor(0x88, 0x88, 0x88),
    'text_muted': RGBColor(0x44, 0x44, 0x44),
    'accent_intro': RGBColor(0x00, 0xD4, 0xAA),
    'accent_rutina': RGBColor(0x00, 0xD4, 0xFF),
    'accent_vacaciones': RGBColor(0xF5, 0xA6, 0x23),
    'accent_conclusiones': RGBColor(0xFF, 0x6B, 0x6B),
    'accent_recap': RGBColor(0xB0, 0xB0, 0xB0),
}

# Colores modo CLARO
COLORS_LIGHT = {
    'bg_primary': RGBColor(0xFA, 0xFA, 0xFA),
    'bg_card': RGBColor(0xF0, 0xF0, 0xF0),
    'text_primary': RGBColor(0x1A, 0x1A, 0x1A),
    'text_secondary': RGBColor(0x55, 0x55, 0x55),
    'text_muted': RGBColor(0xAA, 0xAA, 0xAA),
    'accent_intro': RGBColor(0x00, 0x9E, 0x7A),      # Verde más oscuro para contraste
    'accent_rutina': RGBColor(0x00, 0x9E, 0xCC),     # Cyan más oscuro
    'accent_vacaciones': RGBColor(0xD4, 0x85, 0x00), # Naranja más oscuro
    'accent_conclusiones': RGBColor(0xE0, 0x4A, 0x4A), # Coral más oscuro
    'accent_recap': RGBColor(0x70, 0x70, 0x70),
}

# Por defecto, modo oscuro
COLORS = COLORS_DARK

# Secciones del menú (colores se asignan dinámicamente)
SECTIONS_DATA = [
    {'id': 'intro', 'label': 'Introducción', 'color_key': 'accent_intro'},
    {'id': 'rutina', 'label': 'Mi rutina', 'color_key': 'accent_rutina'},
    {'id': 'vacaciones', 'label': 'Vacaciones', 'color_key': 'accent_vacaciones'},
    {'id': 'conclusiones', 'label': 'Conclusiones', 'color_key': 'accent_conclusiones'},
    {'id': 'recap', 'label': 'Recap', 'color_key': 'accent_recap'},
]

def get_sections(colors):
    """Genera las secciones con los colores del tema."""
    return [
        {'id': s['id'], 'label': s['label'], 'color': colors[s['color_key']]}
        for s in SECTIONS_DATA
    ]

# Contenido de las diapositivas (sin colores hardcodeados)
SLIDES_DATA = [
    # INTRO
    {'section': 'intro', 'type': 'title', 'eyebrow': 'Exposición oral · 2º ESO',
     'title': 'Cómo es mi vida', 'accent_word': 'mi vida', 'subtitle': 'Adri, profesor de Lengua'},
    {'section': 'intro', 'type': 'big_number', 'eyebrow': '01 Datos básicos',
     'number': '42', 'label': 'años'},
    {'section': 'intro', 'type': 'emoji_row', 'eyebrow': '01 Datos básicos',
     'emoji': '👨‍👩‍👧', 'title': 'Mi familia', 'accent_word': 'familia', 'subtitle': 'Pareja + hija de 7 años'},
    {'section': 'intro', 'type': 'emoji_row', 'eyebrow': '01 Datos básicos',
     'emoji': '🏠', 'title': 'Guntín', 'accent_word': 'Guntín', 'subtitle': 'Casa de pueblo'},
    {'section': 'intro', 'type': 'emoji_row', 'eyebrow': '01 Datos básicos',
     'emoji': '🚗', 'title': '35 minutos', 'accent_word': '35', 'subtitle': 'Al instituto cada día'},
    # RUTINA
    {'section': 'rutina', 'type': 'section_title', 'eyebrow': '02 Mi rutina',
     'title': 'Un día normal', 'accent_word': 'normal'},
    {'section': 'rutina', 'type': 'emoji_row', 'eyebrow': '02 Mi rutina',
     'emoji': '⏰', 'title': 'Madrugar', 'accent_word': 'Madrugar', 'subtitle': 'Para llegar a tiempo al instituto'},
    {'section': 'rutina', 'type': 'comparison', 'eyebrow': '02 Mi rutina',
     'title': 'El mediodía', 'accent_word': 'mediodía',
     'left_emoji': '🍽️', 'left_label': 'Comer', 'right_emoji': '🏃', 'right_label': 'Deporte'},
    {'section': 'rutina', 'type': 'emoji_row', 'eyebrow': '02 Mi rutina',
     'emoji': '🌅', 'title': 'Las tardes', 'accent_word': 'tardes', 'subtitle': 'Hija · Clases · Descanso'},
    # VACACIONES
    {'section': 'vacaciones', 'type': 'section_title', 'eyebrow': '03 Vacaciones',
     'title': 'Tiempo de descanso', 'accent_word': 'descanso'},
    {'section': 'vacaciones', 'type': 'emoji_row', 'eyebrow': '03 Vacaciones · Invierno',
     'emoji': '🏝️', 'title': 'Canarias', 'accent_word': 'Canarias', 'subtitle': 'Con mis padres · Energía para el 2º trimestre'},
    {'section': 'vacaciones', 'type': 'emoji_row', 'eyebrow': '03 Vacaciones · Verano',
     'emoji': '🏖️', 'title': 'Playa gallega', 'accent_word': 'gallega', 'subtitle': 'Fiestas familiares + viajes'},
    # CONCLUSIONES
    {'section': 'conclusiones', 'type': 'cards', 'eyebrow': '04 Conclusiones',
     'title': 'Lo que pienso', 'accent_word': 'pienso',
     'card1_title': '✓ Me gusta', 'card1_items': ['Mi familia', 'Galicia', 'Hacer deporte'],
     'card2_title': '✗ Cambiaría', 'card2_items': ['La distancia', 'Casa propia', 'El tiempo gallego']},
    # RECAP
    {'section': 'recap', 'type': 'recap', 'eyebrow': 'Recap',
     'title': 'Mi vida en 4 puntos', 'accent_word': 'vida',
     'recap_items': [
         {'emoji': '👨‍👩‍👧', 'text': '42 años · Guntín\nPareja + hija', 'color_key': 'accent_intro'},
         {'emoji': '⏰', 'text': 'Madrugar · Deporte\nTardes en familia', 'color_key': 'accent_rutina'},
         {'emoji': '🏖️', 'text': 'Canarias · Playa\nFiestas familiares', 'color_key': 'accent_vacaciones'},
         {'emoji': '💭', 'text': 'Me gusta: familia\nCambiaría: distancia', 'color_key': 'accent_conclusiones'},
     ]},
]

def get_slides(colors):
    """Genera las diapositivas con los colores del tema."""
    slides = []
    for slide_data in SLIDES_DATA:
        slide = slide_data.copy()
        # Convertir recap_items a items con colores resueltos
        if 'recap_items' in slide:
            slide['items'] = [
                {'emoji': item['emoji'], 'text': item['text'], 'color': colors[item['color_key']]}
                for item in slide['recap_items']
            ]
            del slide['recap_items']
        slides.append(slide)
    return slides


def set_slide_background(slide, color):
    """Establece el fondo de la diapositiva."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_nav_menu(slide, active_section, sections, colors):
    """Añade el menú de navegación lateral (compacto)."""
    nav_left = Inches(0.15)
    nav_top = Inches(2.2)
    dot_size = Inches(0.08)
    spacing = Inches(0.28)

    for i, section in enumerate(sections):
        is_active = section['id'] == active_section
        y_pos = nav_top + (i * spacing)

        # Punto/círculo pequeño
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, nav_left, y_pos, dot_size, dot_size)
        dot.fill.solid()
        if is_active:
            dot.fill.fore_color.rgb = section['color']
            dot.line.color.rgb = section['color']
            dot.line.width = Pt(1)
        else:
            dot.fill.fore_color.rgb = colors['text_muted']
            dot.line.fill.background()

        # Label pequeño (solo visible si está activo)
        if is_active:
            label = slide.shapes.add_textbox(nav_left + Inches(0.15), y_pos - Inches(0.03), Inches(0.9), Inches(0.2))
            tf = label.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = section['label'].upper()
            p.font.size = Pt(6)
            p.font.bold = True
            p.font.color.rgb = colors['text_secondary']
            p.font.name = FONT_NAME


def add_counter(slide, current, total, color):
    """Añade el contador de diapositivas con color de sección."""
    counter = slide.shapes.add_textbox(Inches(9.2), Inches(0.2), Inches(0.6), Inches(0.25))
    tf = counter.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT

    # Número actual (color de la sección)
    run1 = p.add_run()
    run1.text = str(current)
    run1.font.size = Pt(10)
    run1.font.bold = True
    run1.font.color.rgb = color
    run1.font.name = 'Arial'

    # Separador y total
    run2 = p.add_run()
    run2.text = f'/{total}'
    run2.font.size = Pt(10)
    run2.font.color.rgb = COLORS['text_secondary']
    run2.font.name = 'Arial'


def add_eyebrow(slide, text, color):
    """Añade el eyebrow (etiqueta superior)."""
    # Centrado considerando menú lateral (~1" izquierda): área útil 1"-10", centro en 5.5"
    eyebrow = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(7.5), Inches(0.3))
    tf = eyebrow.text_frame
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    return eyebrow


def add_title_with_accent(slide, title, accent_word, color, font_size=Pt(32)):
    """Añade título con palabra destacada en color."""
    # Centrado considerando menú lateral
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    # Dividir el título en partes
    if accent_word in title:
        parts = title.split(accent_word)
        if parts[0]:
            run = p.add_run()
            run.text = parts[0]
            run.font.size = font_size
            run.font.bold = True
            run.font.color.rgb = COLORS['text_primary']
            run.font.name = FONT_NAME

        run_accent = p.add_run()
        run_accent.text = accent_word
        run_accent.font.size = font_size
        run_accent.font.bold = True
        run_accent.font.color.rgb = color
        run_accent.font.name = FONT_NAME

        if len(parts) > 1 and parts[1]:
            run = p.add_run()
            run.text = parts[1]
            run.font.size = font_size
            run.font.bold = True
            run.font.color.rgb = COLORS['text_primary']
            run.font.name = FONT_NAME
    else:
        run = p.add_run()
        run.text = title
        run.font.size = font_size
        run.font.bold = True
        run.font.color.rgb = COLORS['text_primary']
        run.font.name = FONT_NAME

    return title_box


def add_subtitle(slide, text, top=Inches(2.6)):
    """Añade subtítulo."""
    # Centrado considerando menú lateral
    subtitle = slide.shapes.add_textbox(Inches(1.5), top, Inches(7.5), Inches(0.4))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text_secondary']
    p.font.name = FONT_NAME
    return subtitle


def add_big_number(slide, number, label, color):
    """Añade número grande con etiqueta (centrado en slide)."""
    # Número (centrado considerando menú lateral)
    num_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(7.5), Inches(1.2))
    tf = num_box.text_frame
    tf.anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = number
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT_NAME

    # Label (centrado)
    label_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(7.5), Inches(0.35))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label.upper()
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['text_secondary']
    p.font.name = FONT_NAME


def add_emoji_box(slide, emoji, left, top, size, border_color):
    """Añade caja con emoji centrado manualmente."""
    # Caja de fondo
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size)
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['bg_card']
    box.line.color.rgb = border_color
    box.line.width = Pt(1.5)

    # Emoji centrado manualmente (offset ~30% desde arriba para centrar visualmente)
    emoji_height = Inches(0.5)
    emoji_top = top + (size - emoji_height) / 2
    emoji_box = slide.shapes.add_textbox(left, emoji_top, size, emoji_height)
    tf = emoji_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = emoji
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)

    return box


def add_emoji_row(slide, emoji, title, accent_word, subtitle, color):
    """Añade fila con emoji a la izquierda y contenido a la derecha (centrado en slide)."""
    # Posición vertical centrada (slide altura 5.625")
    box_size = Inches(0.85)
    center_y = Inches(2.45)

    # Emoji box (desplazado a la derecha para compensar menú)
    add_emoji_box(slide, emoji, Inches(3.4), center_y, box_size, color)

    # Título alineado con el centro del emoji box
    title_y = center_y + Inches(0.1)
    title_box = slide.shapes.add_textbox(Inches(4.4), title_y, Inches(4.5), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    if accent_word in title:
        parts = title.split(accent_word)
        if parts[0]:
            run = p.add_run()
            run.text = parts[0]
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = COLORS['text_primary']
            run.font.name = FONT_NAME

        run_accent = p.add_run()
        run_accent.text = accent_word
        run_accent.font.size = Pt(22)
        run_accent.font.bold = True
        run_accent.font.color.rgb = color
        run_accent.font.name = FONT_NAME

        if len(parts) > 1 and parts[1]:
            run = p.add_run()
            run.text = parts[1]
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = COLORS['text_primary']
            run.font.name = FONT_NAME
    else:
        run = p.add_run()
        run.text = title
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = COLORS['text_primary']
        run.font.name = FONT_NAME

    # Subtítulo justo debajo del título
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(4.4), title_y + Inches(0.4), Inches(4.5), Inches(0.3))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['text_secondary']
        p.font.name = FONT_NAME


def add_comparison(slide, title, accent_word, left_emoji, left_label, right_emoji, right_label, color):
    """Añade slide de comparación (A → B) centrado."""
    # Título
    add_title_with_accent(slide, title, accent_word, color, Pt(24))

    # Posición vertical centrada - elementos centrados en área útil (centro ~5.25")
    box_top = Inches(2.45)
    box_size = Inches(0.8)

    # Centrar grupo: box(0.8) + gap(0.25) + arrow(0.4) + gap(0.25) + box(0.8) = 2.5" total
    # Centro en 5.25", inicio en 5.25 - 1.25 = 4.0"
    left_x = Inches(4.0)
    arrow_x = Inches(5.05)
    right_x = Inches(5.7)

    # Elemento izquierdo (dimmed)
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, box_top, box_size, box_size)
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLORS['bg_card']
    left_box.line.color.rgb = COLORS['text_muted']
    left_box.line.width = Pt(1)

    # Emoji izquierdo centrado manualmente
    emoji_height = Inches(0.4)
    emoji_offset = (box_size - emoji_height) / 2
    left_emoji_box = slide.shapes.add_textbox(left_x, box_top + emoji_offset, box_size, emoji_height)
    tf = left_emoji_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_emoji
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)

    left_label_box = slide.shapes.add_textbox(left_x - Inches(0.15), box_top + box_size + Inches(0.06), Inches(1.1), Inches(0.22))
    tf = left_label_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_label
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS['text_muted']
    p.font.name = FONT_NAME

    # Flecha centrada
    arrow = slide.shapes.add_textbox(arrow_x, box_top + emoji_offset, Inches(0.4), emoji_height)
    tf = arrow.text_frame
    p = tf.paragraphs[0]
    p.text = '→'
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['text_muted']

    # Elemento derecho (highlighted)
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, box_top, box_size, box_size)
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(0x00, 0x1A, 0x22)
    right_box.line.color.rgb = color
    right_box.line.width = Pt(1.5)

    # Emoji derecho centrado manualmente
    right_emoji_box = slide.shapes.add_textbox(right_x, box_top + emoji_offset, box_size, emoji_height)
    tf = right_emoji_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_emoji
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)

    right_label_box = slide.shapes.add_textbox(right_x - Inches(0.15), box_top + box_size + Inches(0.06), Inches(1.1), Inches(0.22))
    tf = right_label_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_label
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT_NAME


def add_cards(slide, title, accent_word, card1_title, card1_items, card2_title, card2_items, color):
    """Añade slide con dos tarjetas compactas y centradas."""
    # Título
    add_title_with_accent(slide, title, accent_word, color, Pt(22))

    # Tarjetas más compactas ajustadas al contenido (centradas)
    card_top = Inches(2.35)
    card_width = Inches(1.9)
    card_height = Inches(1.15)

    # Card 1 (Me gusta - verde) - desplazado a la derecha
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.3), card_top, card_width, card_height)
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLORS['bg_card']
    card1.line.color.rgb = RGBColor(0x22, 0x22, 0x22)
    card1.line.width = Pt(1)

    # Título card 1
    c1_title = slide.shapes.add_textbox(Inches(3.4), card_top + Inches(0.12), card_width - Inches(0.2), Inches(0.22))
    tf = c1_title.text_frame
    p = tf.paragraphs[0]
    p.text = card1_title
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_intro']
    p.font.name = FONT_NAME

    # Items card 1 (espaciado ajustado)
    c1_items = slide.shapes.add_textbox(Inches(3.4), card_top + Inches(0.34), card_width - Inches(0.2), Inches(0.75))
    tf = c1_items.text_frame
    for item in card1_items:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['text_secondary']
        p.font.name = FONT_NAME
        p.space_after = Pt(2)

    # Card 2 (Cambiaría - rojo)
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), card_top, card_width, card_height)
    card2.fill.solid()
    card2.fill.fore_color.rgb = COLORS['bg_card']
    card2.line.color.rgb = RGBColor(0x22, 0x22, 0x22)
    card2.line.width = Pt(1)

    # Título card 2
    c2_title = slide.shapes.add_textbox(Inches(5.5), card_top + Inches(0.12), card_width - Inches(0.2), Inches(0.22))
    tf = c2_title.text_frame
    p = tf.paragraphs[0]
    p.text = card2_title
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_conclusiones']
    p.font.name = FONT_NAME

    # Items card 2 (espaciado ajustado)
    c2_items = slide.shapes.add_textbox(Inches(5.5), card_top + Inches(0.34), card_width - Inches(0.2), Inches(0.75))
    tf = c2_items.text_frame
    for item in card2_items:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(9)
        p.font.color.rgb = COLORS['text_secondary']
        p.font.name = FONT_NAME
        p.space_after = Pt(2)


def add_recap(slide, title, accent_word, items):
    """Añade slide de recapitulación con grid centrado."""
    # Título
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    parts = title.split(accent_word)
    if parts[0]:
        run = p.add_run()
        run.text = parts[0]
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = COLORS['text_primary']
        run.font.name = FONT_NAME

    run_accent = p.add_run()
    run_accent.text = accent_word
    run_accent.font.size = Pt(22)
    run_accent.font.bold = True
    run_accent.font.color.rgb = COLORS['text_primary']
    run_accent.font.name = FONT_NAME

    if len(parts) > 1 and parts[1]:
        run = p.add_run()
        run.text = parts[1]
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = COLORS['text_primary']
        run.font.name = FONT_NAME

    # Grid de items - tarjetas más compactas (centrado considerando menú)
    item_width = Inches(1.7)
    item_height = Inches(1.1)
    gap = Inches(0.2)
    total_width = (item_width * 4) + (gap * 3)
    # Centro del área útil (1" a 10") = 5.5", desplazamos 0.5" a la derecha
    start_left = (Inches(10) - total_width) / 2 + Inches(0.5)
    card_top = Inches(2.3)

    for i, item in enumerate(items):
        left = start_left + (i * (item_width + gap))

        # Card (más compacta)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, item_width, item_height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['bg_card']
        card.line.color.rgb = RGBColor(0x22, 0x22, 0x22)
        card.line.width = Pt(1)

        # Emoji centrado manualmente
        emoji_box = slide.shapes.add_textbox(left, card_top + Inches(0.15), item_width, Inches(0.35))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = item['emoji']
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)

        # Texto justo debajo del emoji
        text_box = slide.shapes.add_textbox(left + Inches(0.05), card_top + Inches(0.5), item_width - Inches(0.1), Inches(0.55))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item['text']
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS['text_secondary']
        p.font.name = FONT_NAME


def get_section_color(section_id, colors):
    """Obtiene el color de acento de una sección."""
    color_map = {
        'intro': colors['accent_intro'],
        'rutina': colors['accent_rutina'],
        'vacaciones': colors['accent_vacaciones'],
        'conclusiones': colors['accent_conclusiones'],
        'recap': colors['accent_recap'],
    }
    return color_map.get(section_id, colors['accent_intro'])


def create_presentation(mode='dark'):
    """Crea la presentación completa.

    Args:
        mode: 'dark' o 'light'
    """
    # Seleccionar paleta de colores
    colors = COLORS_DARK if mode == 'dark' else COLORS_LIGHT

    # Actualizar variable global para funciones que la usan
    global COLORS
    COLORS = colors

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # Obtener slides y secciones con colores actualizados
    slides = get_slides(colors)
    sections = get_sections(colors)
    total_slides = len(slides)

    for i, slide_data in enumerate(slides):
        # Crear slide en blanco
        blank_layout = prs.slide_layouts[6]  # Layout en blanco
        slide = prs.slides.add_slide(blank_layout)

        # Fondo según modo
        set_slide_background(slide, colors['bg_primary'])

        # Obtener color de la sección
        section = slide_data['section']
        color = get_section_color(section, colors)

        # Menú de navegación (pasar secciones y colores)
        add_nav_menu(slide, section, sections, colors)

        # Contador (con color de sección)
        add_counter(slide, i + 1, total_slides, color)

        # Eyebrow si existe
        if 'eyebrow' in slide_data:
            add_eyebrow(slide, slide_data['eyebrow'], color)

        # Contenido según tipo
        slide_type = slide_data['type']

        if slide_type == 'title':
            add_title_with_accent(slide, slide_data['title'], slide_data['accent_word'], color)
            add_subtitle(slide, slide_data['subtitle'])

        elif slide_type == 'section_title':
            add_title_with_accent(slide, slide_data['title'], slide_data['accent_word'], color, Pt(48))

        elif slide_type == 'big_number':
            add_big_number(slide, slide_data['number'], slide_data['label'], color)

        elif slide_type == 'emoji_row':
            add_emoji_row(slide, slide_data['emoji'], slide_data['title'],
                         slide_data['accent_word'], slide_data.get('subtitle', ''), color)

        elif slide_type == 'comparison':
            add_comparison(slide, slide_data['title'], slide_data['accent_word'],
                          slide_data['left_emoji'], slide_data['left_label'],
                          slide_data['right_emoji'], slide_data['right_label'], color)

        elif slide_type == 'cards':
            add_cards(slide, slide_data['title'], slide_data['accent_word'],
                     slide_data['card1_title'], slide_data['card1_items'],
                     slide_data['card2_title'], slide_data['card2_items'], color)

        elif slide_type == 'recap':
            add_recap(slide, slide_data['title'], slide_data['accent_word'], slide_data['items'])

    return prs


if __name__ == '__main__':
    base_dir = os.path.dirname(__file__)

    # Generar versión OSCURA
    prs_dark = create_presentation(mode='dark')
    dark_path = os.path.join(base_dir, 'presentacion-mi-vida-oscuro.pptx')
    prs_dark.save(dark_path)
    print(f'✓ Modo oscuro: {dark_path}')

    # Generar versión CLARA
    prs_light = create_presentation(mode='light')
    light_path = os.path.join(base_dir, 'presentacion-mi-vida-claro.pptx')
    prs_light.save(light_path)
    print(f'✓ Modo claro: {light_path}')
